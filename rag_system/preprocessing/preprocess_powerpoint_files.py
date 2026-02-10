from pptx import Presentation
from pptx.slide import Slide
from pptx.table import Table
from pptx.shapes.autoshape import Shape
from typing import List, Dict, Tuple
import re
import time
import io
from .text_utils import clean_text, split_text_into_chunks, normalize_chunk

def process_table(table: Table) -> str:
    """
    Process a PowerPoint table and return its content in a structured text format.
    """
    print(f"Processing table with {len(table.rows)} rows...")
    table_text = []
    for row_idx, row in enumerate(table.rows):
        row_text = []
        for cell in row.cells:
            # Clean cell text and add it to row
            cell_text = clean_text(cell.text)
            if cell_text:
                row_text.append(cell_text)
        if row_text:
            table_text.append(" | ".join(row_text))
        print(f"Processed row {row_idx + 1}/{len(table.rows)}")
    final_output =  "\n".join(table_text)

    return final_output.join("End of table content.\n\n")

def extract_text_from_slide(slide: Slide) -> List[Tuple[str, str]]:
    """
    Extract text from a PowerPoint slide, including tables and shapes.
    Returns a list of tuples (content_type, content) to preserve order.
    """
    content_parts = []
    
    # Process slide title
    if slide.shapes.title:
        title_text = slide.shapes.title.text.strip()
        if title_text:
            content_parts.append(('title', f" {title_text}"))
    
    # Process all shapes in the slide
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            if isinstance(shape, Table):
                table_text = process_table(shape)
                if table_text:
                    content_parts.append(('table', table_text))
            else:
                # Regular text shape
                content_parts.append(('text', shape.text))
    
    return content_parts

def extract_text_from_pptx(file) -> List[Tuple[str, str]]:
    """
    Extract text from a PowerPoint presentation, including all slides and their content.
    Returns a list of tuples (content_type, content) to preserve order.
    """
    try:
        print(f"\nOpening presentation: {file["name"]}")
        prs = Presentation(io.BytesIO(file["content"]))
        all_content = []
        
        print("\nProcessing slides in order...")
        for slide_idx, slide in enumerate(prs.slides):
            # print(f"Processing slide {slide_idx + 1}/{len(prs.slides)}")
            slide_content = extract_text_from_slide(slide)
            if slide_content:
                all_content.extend(slide_content)
        
        print("\nPresentation processing completed")
        return all_content
    except Exception as e:
        print(f"\nError during presentation processing: {str(e)}")
        raise Exception(f"Error extracting text from presentation: {str(e)}")


def preprocess_presentation(file, chunk_size: int = 1000, overlap: int = 200, blob_metadata=None) -> Dict[str, List[str]]:
    """
    Main preprocessing function that handles the entire presentation preprocessing pipeline.
    """
    print(f"\nStarting presentation preprocessing: {file["name"]}")
    
    print("\nStep 1: Extracting text from presentation")
    content_parts = extract_text_from_pptx(file)
    print(f"Extracted {len(content_parts)} content parts from presentation")
    
    # Process each part according to its type
    processed_parts = []
    for content_type, content in content_parts:
        if content_type == 'table':
            # Keep tables as is
            processed_parts.append(f"\nTable Content:\n{content}")
        else:
            # Clean non-table content
            processed_parts.append(clean_text(content))
    
    # Join all parts together
    final_text = '\n'.join(processed_parts)
    # final_text = final_text.replace("JPEG 2000","")
    print(f"Extracted {len(final_text)} characters of processed text")
    print('Processed text preview:', final_text[:500])
    
    print("\nStep 2: Splitting into chunks")
    chunks = split_text_into_chunks(final_text, chunk_size, overlap)
    
    # Normalize chunks
    normalized_chunks = []
    for chunk in chunks:
        normalized_chunk = normalize_chunk(chunk)
        normalized_chunks.append(normalized_chunk)
        print(f"Normalized chunk: {normalized_chunk[:100]}...")
    
    print(f"Chunk length for presentation {file['name']} = {len(normalized_chunks)}")
    
    # Prepare metadata
    metadata = {
        'user_id': 'example_user',
        'client_id': 'example_client',
        'timestamp_in_seconds': time.time(),
        'num_chunks': len(normalized_chunks),
        'total_chars': len(final_text)
    }
    metadata.update(blob_metadata or {})
    metadata["has_been_preprocessed"] = "True"
    
    print("\nPreprocessing completed successfully!")
    return {
        'result': {
            'chunks': normalized_chunks,
            'metadata': metadata
        }
    }

if __name__ == "__main__":
    # Example usage
    try:
        start_time = time.time()
        pptx_path = r"test_data/6_Image compression.pptx"
        print("\n=== Starting Presentation Processing ===")
        
        # Read the file and create the expected dictionary format
        with open(pptx_path, 'rb') as f:
            file_content = f.read()
        
        file_dict = {"name": pptx_path, "content": file_content}
        result = preprocess_presentation(file_dict)
        
        print("\n=== Processing Results ===")
        print(f"Processed presentation: {result['result']['metadata']['user_id']}")
        print(f"Number of chunks: {result['result']['metadata']['num_chunks']}")
        print(f"Total characters: {result['result']['metadata']['total_chars']}")
        print(f"Processing time: {time.time() - start_time:.2f} seconds")
        
        # Print first chunk as example
        if result['result']['chunks']:
            print(f"\nFirst chunk example:\n{result['result']['chunks'][0]}\n")
            
    except Exception as e:
        print(f"\nError processing presentation: {str(e)}")
