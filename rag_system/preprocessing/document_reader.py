import os
import json
import re 
from pathlib import Path
from typing import List, Dict, Any, Union, Optional
import chromadb
from chromadb.utils import embedding_functions
import tiktoken 
from .preprocess_pdf_files import split_text_into_chunks

class DocumentReader:
    def __init__(self, chroma_db_path="./chroma_db2"):
        self.client = chromadb.PersistentClient(path=chroma_db_path)
        # Simplified supported extensions - only the ones we need
        self.supported_extensions = {'.txt', '.pdf', '.docx', '.pptx', '.ppt', '.doc'} 
        self.default_embedding_function = embedding_functions.SentenceTransformerEmbeddingFunction(
            model_name="sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2"
        )

    def info(self):
        """Prints information about the DocumentReader instance."""
        try:
            collection = self.client.get_or_create_collection(
                name="multimodal_data",
                embedding_function=self.default_embedding_function
            )
            count = collection.count()
            print(f"Collection '{collection.name}' has {count} documents.")
            if count > 0:
                results = collection.peek(limit=5)
                metadatas = results.get('metadatas', [])
                documents = results.get('documents', [])
                print("Sample Documents & Metadatas (first 5 documents):")
                for i in range(min(len(documents), len(metadatas))):
                    print(f"  Content (first 100 chars): {documents[i][:100]}...")
                    print(f"  Metadata: {metadatas[i]}")
            else:
                print("No documents in collection to display metadatas.")
        except Exception as e:
            print(f"Error accessing collection info: {e}")

    def _process_text_content(self, file_path: str, content: str, file_type: str, 
                            chunk_size: int = None, overlap: int = None) -> List[Dict[str, Any]]:
        """Simplified method for processing text content."""
        
        # Use default chunking parameters if not provided
        chunk_size = chunk_size or 1500
        overlap = overlap or 200
        
        # Use the existing chunking function
        text_chunks = split_text_into_chunks(content, chunk_size=chunk_size, overlap=overlap)
        
        base_metadata = {
            'file_type': file_type,
            'file_name': os.path.basename(file_path),
            'file_path': file_path,
            'file_size': os.path.getsize(file_path) if os.path.exists(file_path) else 0,
            'language': 'macedonian'
        }
        
        documents = []
        for i, chunk in enumerate(text_chunks):
            processed_chunk = chunk.strip()
            if processed_chunk and len(processed_chunk) > 50:  # Only keep meaningful chunks
                metadata = {
                    **base_metadata,
                    'chunk_index': i,
                    'chunk_length': len(processed_chunk),
                    'num_words': len(processed_chunk.split())
                }
                documents.append({'content': processed_chunk, 'metadata': metadata})
        
        return documents

    def read_text_file(self, file_path: str, chunk_size: int = None, overlap: int = None) -> Union[List[Dict[str, Any]], None]:
        """Read and process text files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Clean up content
            content = re.sub(r'\s+', ' ', content).strip()
            return self._process_text_content(file_path, content, 'txt', chunk_size, overlap)
        except Exception as e:
            print(f"Error reading text file {file_path}: {e}")
            return None

    def _process_pdf_file(self, file_path: str, chunk_size: int = None, overlap: int = None) -> Union[List[Dict[str, Any]], None]:
        """Process PDF files using the existing preprocessing script."""
        try:
            # Import PDF processing function
            import importlib.util
            script_path = os.path.join(os.path.dirname(__file__), 'preprocess_pdf_files.py')
            
            spec = importlib.util.spec_from_file_location("pdf_processor", script_path)
            pdf_module = importlib.util.module_from_spec(spec)
            spec.loader.exec_module(pdf_module)
            
            # Use the PDF processing function
            result = pdf_module.read_pdf_file(file_path, chunk_size or 1500, overlap or 200)
            
            if result:
                # Simplify the metadata structure
                simplified_docs = []
                for doc in result:
                    metadata = {
                        'file_type': 'pdf',
                        'file_name': os.path.basename(file_path),
                        'file_path': file_path,
                        'file_size': os.path.getsize(file_path),
                        'language': 'macedonian',
                        'chunk_index': doc.get('metadata', {}).get('chunk_index', 0),
                        'chunk_length': len(doc.get('content', '')),
                        'num_words': len(doc.get('content', '').split())
                    }
                    
                    # Add any additional metadata from original processing
                    original_metadata = doc.get('metadata', {})
                    if 'page_number' in original_metadata:
                        metadata['page_number'] = original_metadata['page_number']
                    
                    simplified_docs.append({
                        'content': doc.get('content', ''),
                        'metadata': metadata
                    })
                
                return simplified_docs
            return None
        except Exception as e:
            print(f"Error processing PDF {file_path}: {e}")
            return None

    def _process_docx_file(self, file_path: str, chunk_size: int = None, overlap: int = None) -> Union[List[Dict[str, Any]], None]:
        """Process Word document files."""
        try:
            # Import docx processing
            import importlib.util
            script_path = os.path.join(os.path.dirname(__file__), 'preprocess_doc_files.py')
            
            if os.path.exists(script_path):
                spec = importlib.util.spec_from_file_location("doc_processor", script_path)
                doc_module = importlib.util.module_from_spec(spec)
                spec.loader.exec_module(doc_module)
                
                result = doc_module.read_docx_file(file_path, chunk_size or 1500, overlap or 200)
                
                if result:
                    # Simplify metadata
                    simplified_docs = []
                    for doc in result:
                        metadata = {
                            'file_type': 'docx',
                            'file_name': os.path.basename(file_path),
                            'file_path': file_path,
                            'file_size': os.path.getsize(file_path),
                            'language': 'macedonian',
                            'chunk_index': doc.get('metadata', {}).get('chunk_index', 0),
                            'chunk_length': len(doc.get('content', '')),
                            'num_words': len(doc.get('content', '').split())
                        }
                        
                        simplified_docs.append({
                            'content': doc.get('content', ''),
                            'metadata': metadata
                        })
                    
                    return simplified_docs
            else:
                # Fallback: try to extract text manually
                try:
                    from docx import Document
                    doc = Document(file_path)
                    content = ""
                    for paragraph in doc.paragraphs:
                        content += paragraph.text + "\\n"
                    
                    content = re.sub(r'\\s+', ' ', content).strip()
                    return self._process_text_content(file_path, content, 'docx', chunk_size, overlap)
                except ImportError:
                    print(f"python-docx not installed, skipping {file_path}")
                    return None
                    
        except Exception as e:
            print(f"Error processing Word document {file_path}: {e}")
            return None

    def _process_pptx_file(self, file_path: str, chunk_size: int = None, overlap: int = None) -> Union[List[Dict[str, Any]], None]:
        """Process PowerPoint files."""
        try:
            # Import pptx processing
            import importlib.util
            script_path = os.path.join(os.path.dirname(__file__), 'preprocess_powerpoint_files.py')
            
            if os.path.exists(script_path):
                spec = importlib.util.spec_from_file_location("ppt_processor", script_path)
                ppt_module = importlib.util.module_from_spec(spec)
                spec.loader.exec_module(ppt_module)
                
                result = ppt_module.read_pptx_file(file_path, chunk_size or 1500, overlap or 200)
                
                if result:
                    # Simplify metadata
                    simplified_docs = []
                    for doc in result:
                        metadata = {
                            'file_type': 'pptx',
                            'file_name': os.path.basename(file_path),
                            'file_path': file_path,
                            'file_size': os.path.getsize(file_path),
                            'language': 'macedonian',
                            'chunk_index': doc.get('metadata', {}).get('chunk_index', 0),
                            'chunk_length': len(doc.get('content', '')),
                            'num_words': len(doc.get('content', '').split()),
                            'slide_number': doc.get('metadata', {}).get('slide_number', 0)
                        }
                        
                        simplified_docs.append({
                            'content': doc.get('content', ''),
                            'metadata': metadata
                        })
                    
                    return simplified_docs
            else:
                # Fallback: try to extract text manually
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    content = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                content += shape.text + "\\n"
                    
                    content = re.sub(r'\\s+', ' ', content).strip()
                    return self._process_text_content(file_path, content, 'pptx', chunk_size, overlap)
                except ImportError:
                    print(f"python-pptx not installed, skipping {file_path}")
                    return None
                    
        except Exception as e:
            print(f"Error processing PowerPoint {file_path}: {e}")
            return None

    def read_single_document(self, file_path: str, chunk_size: int = 1500, overlap: int = 200) -> Union[List[Dict[str, Any]], None]:
        """Read and process a single document based on its extension."""
        
        if not os.path.exists(file_path):
            print(f"File does not exist: {file_path}")
            return None
        
        file_extension = Path(file_path).suffix.lower()
        
        if file_extension not in self.supported_extensions:
            print(f"Unsupported file type: {file_extension}")
            return None
        
        print(f"Processing {file_extension} file: {os.path.basename(file_path)}")
        
        # Route to appropriate processor
        if file_extension == '.txt':
            return self.read_text_file(file_path, chunk_size, overlap)
        elif file_extension == '.pdf':
            return self._process_pdf_file(file_path, chunk_size, overlap)
        elif file_extension == '.docx' or file_extension == '.doc':
            return self._process_docx_file(file_path, chunk_size, overlap)
        elif file_extension == '.pptx' or file_extension == '.ppt':
            return self._process_pptx_file(file_path, chunk_size, overlap)
        else:
            print(f"No processor available for {file_extension}")
            return None

    def process_directory(self, directory_path: str, collection_name: str = "macedonian_documents", 
                         chunk_size: int = 1500, overlap: int = 200) -> Dict[str, int]:
        """Process all supported documents in a directory and add them to ChromaDB."""
        
        if not os.path.exists(directory_path):
            print(f"Directory does not exist: {directory_path}")
            return {'processed': 0, 'failed': 0}
        
        # Get or create collection
        collection = self.client.get_or_create_collection(
            name=collection_name,
            embedding_function=self.default_embedding_function
        )
        
        processed_count = 0
        failed_count = 0
        
        # Process all supported files
        for root, dirs, files in os.walk(directory_path):
            for file in files:
                file_path = os.path.join(root, file)
                file_extension = Path(file_path).suffix.lower()
                
                if file_extension in self.supported_extensions:
                    try:
                        documents = self.read_single_document(file_path, chunk_size, overlap)
                        
                        if documents:
                            # Add to ChromaDB
                            ids = []
                            texts = []
                            metadatas = []
                            
                            for i, doc in enumerate(documents):
                                doc_id = f"{collection_name}_{os.path.basename(file_path)}_{i}"
                                ids.append(doc_id)
                                texts.append(doc['content'])
                                metadatas.append(doc['metadata'])
                            
                            collection.add(
                                ids=ids,
                                documents=texts,
                                metadatas=metadatas
                            )
                            
                            processed_count += 1
                            print(f"✅ Processed: {file} ({len(documents)} chunks)")
                        else:
                            failed_count += 1
                            print(f"❌ Failed: {file}")
                            
                    except Exception as e:
                        failed_count += 1
                        print(f"❌ Error processing {file}: {e}")
        
        print(f"\\n📊 Processing complete:")
        print(f"  ✅ Processed: {processed_count} files")
        print(f"  ❌ Failed: {failed_count} files")
        print(f"  📚 Collection: {collection_name}")
        
        return {'processed': processed_count, 'failed': failed_count}

    def add_documents_to_collection(self, documents: List[Dict[str, Any]], collection_name: str = "macedonian_documents") -> bool:
        """Add a list of processed documents to a ChromaDB collection."""
        try:
            # Get or create collection
            collection = self.client.get_or_create_collection(
                name=collection_name,
                embedding_function=self.default_embedding_function
            )
            
            if not documents:
                print("No documents to add")
                return False
            
            # Prepare data for ChromaDB
            ids = []
            texts = []
            metadatas = []
            
            for i, doc in enumerate(documents):
                # Create unique ID based on file name and chunk index
                file_name = doc['metadata'].get('file_name', 'unknown')
                chunk_index = doc['metadata'].get('chunk_index', i)
                doc_id = f"{collection_name}_{file_name}_{chunk_index}_{i}"
                
                ids.append(doc_id)
                texts.append(doc['content'])
                metadatas.append(doc['metadata'])
            
            # Add to collection
            collection.add(
                ids=ids,
                documents=texts,
                metadatas=metadatas
            )
            
            print(f"✅ Added {len(documents)} documents to collection '{collection_name}'")
            return True
            
        except Exception as e:
            print(f"❌ Error adding documents to collection: {e}")
            return False

    def query_documents(self, query: str, collection_name: str = "macedonian_documents", 
                       n_results: int = 10) -> List[str]:
        """Query documents from ChromaDB collection."""
        try:
            collection = self.client.get_collection(name=collection_name)
            
            results = collection.query(
                query_texts=[query],
                n_results=n_results
            )
            
            if results['documents'] and results['documents'][0]:
                return results['documents'][0]
            else:
                return []
                
        except Exception as e:
            print(f"Error querying documents: {e}")
            return []

    def get_collection_stats(self, collection_name: str = "macedonian_documents") -> Dict[str, Any]:
        """Get statistics about a collection."""
        try:
            collection = self.client.get_collection(name=collection_name)
            count = collection.count()
            
            if count > 0:
                # Sample some documents to get metadata info
                sample = collection.peek(limit=min(count, 100))
                
                # Analyze file types
                file_types = {}
                languages = {}
                
                for metadata in sample.get('metadatas', []):
                    file_type = metadata.get('file_type', 'unknown')
                    language = metadata.get('language', 'unknown')
                    
                    file_types[file_type] = file_types.get(file_type, 0) + 1
                    languages[language] = languages.get(language, 0) + 1
                
                return {
                    'total_documents': count,
                    'file_types': file_types,
                    'languages': languages,
                    'collection_name': collection_name
                }
            else:
                return {
                    'total_documents': 0,
                    'file_types': {},
                    'languages': {},
                    'collection_name': collection_name
                }
                
        except Exception as e:
            print(f"Error getting collection stats: {e}")
            return {}
